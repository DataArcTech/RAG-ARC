import os
import json
import logging
from dataclasses import dataclass
from typing import List, Optional
from urllib.parse import urlparse

from .base import ParserBase

logger = logging.getLogger(__name__)


@dataclass
class NativeParser(ParserBase):
    """
    Multi-format document parser supporting PDF, DOCX, Excel, PowerPoint, HTML, and images.

    This class provides a unified parsing interface for multiple document formats,
    automatically routing files to appropriate specialized parsing functions based on file extension.
    Supports both local files and remote URLs with configurable output formats.

    Supported file types:
    - DOCX: Native Word document parsing with image extraction
    - Excel: XLSX/CSV parsing with table structure preservation
    - PowerPoint: PPTX slide content and layout extraction
    - HTML: Web page parsing with content extraction
    - Images: JPG/PNG (requires DotsOCR integration)
    - PDF: (requires DotsOCR integration)
    """

    def parse_file(
        self,
        input_path: str,
        output_dir: Optional[str] = None,
        **kwargs
    ) -> List[dict]:
        """Parse a file of any supported type"""

        # Set up output directory
        output_dir = output_dir or getattr(self.config, 'output_dir', 'output')
        output_dir = os.path.abspath(output_dir)
        os.makedirs(output_dir, exist_ok=True)

        # Handle URLs
        if self._is_url(input_path):
            if self._is_html_url(input_path):
                return self._parse_html_url(input_path, output_dir, **kwargs)
            else:
                raise ValueError(f"Only HTML URLs are supported, got: {input_path}")

        # Validate file exists
        if not os.path.isfile(input_path):
            raise FileNotFoundError(f"File not found: {input_path}")

        # Extract file extension and validate
        _, file_ext = os.path.splitext(input_path)
        file_ext = file_ext.lower()

        if file_ext not in self.get_supported_extensions():
            supported = ', '.join(self.get_supported_extensions())
            raise ValueError(f"Unsupported file type '{file_ext}'. Supported types: {supported}")

        # Route to appropriate parser method
        try:
            if file_ext == '.docx':
                return self._parse_docx(input_path, output_dir, **kwargs)
            elif file_ext in ['.xlsx', '.xls', '.csv']:
                return self._parse_excel(input_path, output_dir, **kwargs)
            elif file_ext == '.pptx':
                return self._parse_ppt(input_path, output_dir, **kwargs)
            elif file_ext == '.html':
                return self._parse_html_file(input_path, output_dir, **kwargs)
            else:
                raise ValueError(f"File type '{file_ext}' is listed as supported but no handler exists")

        except Exception as e:
            logger.error(f"Failed to parse {input_path}: {str(e)}")
            raise RuntimeError(f"Failed to parse {input_path}: {str(e)}")

    def get_supported_extensions(self) -> List[str]:
        """Get all supported file extensions"""
        return ['.docx', '.xlsx', '.xls', '.csv', '.pptx', '.html']

    def _is_url(self, path: str) -> bool:
        """Check if path is a URL"""
        try:
            result = urlparse(path)
            return bool(result.scheme and result.netloc)
        except:
            return False

    def _is_html_url(self, url: str) -> bool:
        """Check if URL points to HTML content"""
        return url.lower().endswith('.html') or not any(
            url.lower().endswith(ext) for ext in self.get_supported_extensions()
        )

    # ==================== PRIVATE PARSING METHODS ====================

    def _parse_docx(self, input_path: str, output_dir: str, **kwargs) -> List[dict]:
        """Parse DOCX file and return structured results"""
        try:
            from docx import Document

            filename = os.path.splitext(os.path.basename(input_path))[0]
            save_dir = os.path.join(output_dir, filename)
            os.makedirs(save_dir, exist_ok=True)

            print(f"Parsing DOCX: {filename}")

            # Parse DOCX content
            doc = Document(input_path)

            # Extract text content
            full_text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    full_text.append(paragraph.text)

            # Extract tables
            tables_data = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    tables_data.append(table_data)

            # Save results
            content = {
                'text': '\n'.join(full_text),
                'tables': tables_data,
                'metadata': {
                    'filename': filename,
                    'paragraphs_count': len([p for p in doc.paragraphs if p.text.strip()]),
                    'tables_count': len(tables_data)
                }
            }

            # Save as JSON
            json_path = os.path.join(save_dir, f"{filename}.json")
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(content, f, ensure_ascii=False, indent=2)

            # Save as Markdown
            md_content = self._convert_docx_to_markdown(content)
            md_path = os.path.join(save_dir, f"{filename}.md")
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(md_content)

            result = {
                'file_path': input_path,
                'page_no': 0,
                'content_type': 'docx',
                'output_paths': {
                    'json': json_path,
                    'markdown': md_path
                },
                'metadata': content['metadata']
            }

            return [result]

        except Exception as e:
            logger.error(f"DOCX parsing failed: {str(e)}")
            raise

    def _parse_excel(self, input_path: str, output_dir: str, **kwargs) -> List[dict]:
        """Parse Excel file and return structured results"""
        try:
            import pandas as pd
            import chardet

            filename = os.path.splitext(os.path.basename(input_path))[0]
            save_dir = os.path.join(output_dir, filename)
            os.makedirs(save_dir, exist_ok=True)

            print(f"Parsing Excel: {filename}")

            # Read all sheets
            if input_path.lower().endswith('.csv'):
                # Detect encoding for CSV
                with open(input_path, 'rb') as f:
                    raw = f.read(10000)
                    result = chardet.detect(raw)
                    encoding = result['encoding'] or "utf-8"
                sheets_data = {'Sheet1': pd.read_csv(input_path, encoding=encoding)}
            else:
                sheets_data = pd.read_excel(input_path, sheet_name=None)

            all_content = []
            sheet_results = []

            for sheet_name, df in sheets_data.items():
                # Convert to structured data
                sheet_content = {
                    'sheet_name': sheet_name,
                    'data': df.to_dict('records'),
                    'columns': df.columns.tolist(),
                    'shape': df.shape,
                    'metadata': {
                        'rows': len(df),
                        'columns': len(df.columns),
                        'empty_cells': int(df.isnull().sum().sum())
                    }
                }
                all_content.append(sheet_content)

                # Save individual sheet as CSV
                csv_path = os.path.join(save_dir, f"{filename}_{sheet_name}.csv")
                df.to_csv(csv_path, index=False, encoding='utf-8')

                # Save as JSON
                json_path = os.path.join(save_dir, f"{filename}_{sheet_name}.json")
                with open(json_path, 'w', encoding='utf-8') as f:
                    json.dump(sheet_content, f, ensure_ascii=False, indent=2)

                sheet_results.append({
                    'file_path': input_path,
                    'page_no': len(sheet_results),
                    'content_type': 'excel_sheet',
                    'sheet_name': sheet_name,
                    'output_paths': {
                        'json': json_path,
                        'csv': csv_path
                    },
                    'metadata': sheet_content['metadata']
                })

            # Save combined results
            combined_json = os.path.join(save_dir, f"{filename}_combined.json")
            with open(combined_json, 'w', encoding='utf-8') as f:
                json.dump(all_content, f, ensure_ascii=False, indent=2)

            return sheet_results

        except Exception as e:
            logger.error(f"Excel parsing failed: {str(e)}")
            raise

    def _parse_ppt(self, input_path: str, output_dir: str, **kwargs) -> List[dict]:
        """Parse PowerPoint file and return structured results"""
        try:
            from pptx import Presentation

            filename = os.path.splitext(os.path.basename(input_path))[0]
            save_dir = os.path.join(output_dir, filename)
            os.makedirs(save_dir, exist_ok=True)

            print(f"Parsing PowerPoint: {filename}")

            prs = Presentation(input_path)
            slides_data = []
            results = []

            for i, slide in enumerate(prs.slides):
                slide_content = {
                    'slide_number': i + 1,
                    'title': '',
                    'text_content': [],
                    'notes': ''
                }

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        if not slide_content['title'] and hasattr(shape, 'text_frame'):
                            slide_content['title'] = shape.text.strip()
                        else:
                            slide_content['text_content'].append(shape.text.strip())

                # Extract notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    slide_content['notes'] = slide.notes_slide.notes_text_frame.text.strip()

                slides_data.append(slide_content)

                # Save individual slide
                slide_json = os.path.join(save_dir, f"{filename}_slide_{i+1}.json")
                with open(slide_json, 'w', encoding='utf-8') as f:
                    json.dump(slide_content, f, ensure_ascii=False, indent=2)

                # Convert to markdown
                md_content = self._convert_slide_to_markdown(slide_content)
                slide_md = os.path.join(save_dir, f"{filename}_slide_{i+1}.md")
                with open(slide_md, 'w', encoding='utf-8') as f:
                    f.write(md_content)

                results.append({
                    'file_path': input_path,
                    'page_no': i,
                    'content_type': 'ppt_slide',
                    'slide_number': i + 1,
                    'output_paths': {
                        'json': slide_json,
                        'markdown': slide_md
                    },
                    'metadata': {
                        'title': slide_content['title'],
                        'text_blocks': len(slide_content['text_content']),
                        'has_notes': bool(slide_content['notes'])
                    }
                })

            # Save combined presentation
            combined_json = os.path.join(save_dir, f"{filename}_combined.json")
            with open(combined_json, 'w', encoding='utf-8') as f:
                json.dump(slides_data, f, ensure_ascii=False, indent=2)

            return results

        except Exception as e:
            logger.error(f"PowerPoint parsing failed: {str(e)}")
            raise

    def _parse_html_file(self, input_path: str, output_dir: str, **kwargs) -> List[dict]:
        """Parse HTML file and return structured results"""
        with open(input_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        filename = os.path.splitext(os.path.basename(input_path))[0]
        return self._parse_html_content(html_content, input_path, filename, output_dir)

    def _parse_html_url(self, url: str, output_dir: str, **kwargs) -> List[dict]:
        """Parse HTML URL and return structured results"""
        try:
            import requests

            print(f"Fetching HTML from URL: {url}")
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            html_content = response.text

            filename = urlparse(url).path.split('/')[-1] or 'webpage'
            filename = os.path.splitext(filename)[0] or 'webpage'

            return self._parse_html_content(html_content, url, filename, output_dir)

        except Exception as e:
            logger.error(f"HTML URL parsing failed: {str(e)}")
            raise

    def _parse_html_content(self, html_content: str, source_path: str, filename: str, output_dir: str) -> List[dict]:
        """Parse HTML content and return structured results"""
        try:
            from bs4 import BeautifulSoup

            save_dir = os.path.join(output_dir, filename)
            os.makedirs(save_dir, exist_ok=True)

            print(f"Parsing HTML: {filename}")

            # Parse HTML
            soup = BeautifulSoup(html_content, 'html.parser')

            # Extract structured content
            content = {
                'title': soup.title.string if soup.title else '',
                'headings': [],
                'paragraphs': [],
                'links': [],
                'images': [],
                'tables': [],
                'metadata': {
                    'url': source_path if source_path.startswith('http') else '',
                    'filename': filename
                }
            }

            # Extract headings
            for tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                for heading in soup.find_all(tag):
                    if heading.get_text().strip():
                        content['headings'].append({
                            'level': tag,
                            'text': heading.get_text().strip()
                        })

            # Extract paragraphs
            for p in soup.find_all('p'):
                text = p.get_text().strip()
                if text:
                    content['paragraphs'].append(text)

            # Extract links
            for a in soup.find_all('a', href=True):
                content['links'].append({
                    'text': a.get_text().strip(),
                    'href': a['href']
                })

            # Extract images
            for img in soup.find_all('img'):
                content['images'].append({
                    'src': img.get('src', ''),
                    'alt': img.get('alt', ''),
                    'title': img.get('title', '')
                })

            # Extract tables
            for table in soup.find_all('table'):
                table_data = []
                for row in table.find_all('tr'):
                    row_data = [cell.get_text().strip() for cell in row.find_all(['td', 'th'])]
                    if row_data:
                        table_data.append(row_data)
                if table_data:
                    content['tables'].append(table_data)

            # Save results
            json_path = os.path.join(save_dir, f"{filename}.json")
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(content, f, ensure_ascii=False, indent=2)

            # Convert to markdown
            md_content = self._convert_html_to_markdown(content)
            md_path = os.path.join(save_dir, f"{filename}.md")
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(md_content)

            result = {
                'file_path': source_path,
                'page_no': 0,
                'content_type': 'html',
                'output_paths': {
                    'json': json_path,
                    'markdown': md_path
                },
                'metadata': {
                    'title': content['title'],
                    'headings_count': len(content['headings']),
                    'paragraphs_count': len(content['paragraphs']),
                    'links_count': len(content['links']),
                    'images_count': len(content['images']),
                    'tables_count': len(content['tables'])
                }
            }

            return [result]

        except Exception as e:
            logger.error(f"HTML parsing failed: {str(e)}")
            raise

    # ==================== PRIVATE UTILITY METHODS ====================

    def _convert_docx_to_markdown(self, content: dict) -> str:
        """Convert DOCX content to Markdown"""
        md_lines = []

        # Add text content
        if content['text']:
            md_lines.append(content['text'])
            md_lines.append('')

        # Add tables
        for i, table in enumerate(content['tables']):
            md_lines.append(f"## Table {i+1}")
            md_lines.append('')

            if table and len(table) > 0:
                # Header row
                header = '| ' + ' | '.join(table[0]) + ' |'
                separator = '| ' + ' | '.join(['---'] * len(table[0])) + ' |'
                md_lines.append(header)
                md_lines.append(separator)

                # Data rows
                for row in table[1:]:
                    row_md = '| ' + ' | '.join(row) + ' |'
                    md_lines.append(row_md)

            md_lines.append('')

        return '\n'.join(md_lines)

    def _convert_slide_to_markdown(self, slide_content: dict) -> str:
        """Convert slide content to Markdown"""
        md_lines = []

        if slide_content['title']:
            md_lines.append(f"# {slide_content['title']}")
            md_lines.append('')

        for text in slide_content['text_content']:
            md_lines.append(text)
            md_lines.append('')

        if slide_content['notes']:
            md_lines.append("## Notes")
            md_lines.append(slide_content['notes'])

        return '\n'.join(md_lines)

    def _convert_html_to_markdown(self, content: dict) -> str:
        """Convert HTML content to Markdown"""
        md_lines = []

        if content['title']:
            md_lines.append(f"# {content['title']}")
            md_lines.append('')

        # Add headings and paragraphs in order they appear
        for heading in content['headings']:
            level = int(heading['level'][1])
            md_lines.append('#' * level + f" {heading['text']}")
            md_lines.append('')

        for paragraph in content['paragraphs']:
            md_lines.append(paragraph)
            md_lines.append('')

        # Add tables
        for i, table in enumerate(content['tables']):
            md_lines.append(f"## Table {i+1}")
            md_lines.append('')

            if table and len(table) > 0:
                # Header row
                header = '| ' + ' | '.join(table[0]) + ' |'
                separator = '| ' + ' | '.join(['---'] * len(table[0])) + ' |'
                md_lines.append(header)
                md_lines.append(separator)

                # Data rows
                for row in table[1:]:
                    row_md = '| ' + ' | '.join(row) + ' |'
                    md_lines.append(row_md)

            md_lines.append('')

        return '\n'.join(md_lines)